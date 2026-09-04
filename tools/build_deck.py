#!/usr/bin/env python3
"""Build a .pptx (with speaker notes in the notes pane) from a deck markdown file.

Usage:
    python tools/build_deck.py docs/coe-sessions/decks/<deck>.deck.md
    python tools/build_deck.py <deck>.deck.md -o /some/path/out.pptx

Deck markdown format
--------------------
    # <deck title>

    <!-- meta
    session: 01
    ppt_date_time: 2026-09-11 14:00 IST
    presenter: Ramjee
    subtitle: optional subtitle for the title slide
    -->

    ## <slide title>
    [TIME: 1:30]
    - bullet, may use **bold** and `code`
      - sub-bullet (two-space indent)
    NOTES:
    Everything after NOTES: until the next '## ' becomes the speaker notes.

Diagrams
--------
Instead of bullets, a slide may declare a diagram. Items are `Label | detail`.

    DIAGRAM: ladder      ascending steps - progression, maturity levels
    DIAGRAM: flow        left-to-right boxes with arrows - a process
    DIAGRAM: compare     side-by-side columns - before/after, with/without
    DIAGRAM: cycle       boxes in a loop - repeating process

Every slide is drawn from a blank layout with explicit geometry, so nothing
inherits or overlaps.

Install once:  pip install python-pptx
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError:  # pragma: no cover
    sys.exit("python-pptx is not installed. Run: pip install python-pptx")

# ---------------------------------------------------------------- palette ---
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
ACCENT = RGBColor(0x0A, 0x6E, 0xD1)
ACCENT_DK = RGBColor(0x05, 0x4A, 0x8F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
# Progressive fills for ladder/flow steps (light -> saturated)
TINTS = [
    RGBColor(0xE8, 0xF1, 0xFB),
    RGBColor(0xC7, 0xDF, 0xF6),
    RGBColor(0x8F, 0xC0, 0xEC),
    RGBColor(0x4E, 0x9B, 0xDE),
    RGBColor(0x0A, 0x6E, 0xD1),
]

# 16:9 canvas
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
# Content region, below the title band
MARGIN = Inches(0.75)
BODY_TOP = Inches(1.75)
BODY_W = Inches(11.83)
BODY_H = Inches(5.15)

FONT = "Segoe UI"
MONO = "Consolas"

INLINE = re.compile(r"(\*\*.+?\*\*|`[^`]+?`)")


# ------------------------------------------------------------------ parse ---
def parse_deck(text: str) -> tuple[dict, list[dict]]:
    meta: dict[str, str] = {}
    if m := re.search(r"<!--\s*meta(.*?)-->", text, re.S):
        for line in m.group(1).strip().splitlines():
            if ":" in line:
                key, _, value = line.partition(":")
                meta[key.strip()] = value.strip()
        text = text[: m.start()] + text[m.end() :]

    if m := re.search(r"^#\s+(.+)$", text, re.M):
        meta.setdefault("title", m.group(1).strip())

    slides: list[dict] = []
    chunks = re.split(r"^##\s+(.+)$", text, flags=re.M)[1:]
    for title, body in zip(chunks[0::2], chunks[1::2]):
        content, _, notes = body.partition("NOTES:")

        time = ""
        if m := re.search(r"^\[TIME:\s*(.+?)\]\s*$", content, re.M):
            time = m.group(1).strip()
            content = content[: m.start()] + content[m.end() :]

        diagram = ""
        if m := re.search(r"^DIAGRAM:\s*(\w+)\s*$", content, re.M):
            diagram = m.group(1).strip().lower()
            content = content[: m.start()] + content[m.end() :]

        caption = ""
        if m := re.search(r"^CAPTION:\s*(.+?)\s*$", content, re.M):
            caption = m.group(1).strip()
            content = content[: m.start()] + content[m.end() :]

        items: list[tuple[int, str]] = []
        for raw in content.splitlines():
            if not raw.strip().startswith(("-", "*")):
                continue
            indent = len(raw) - len(raw.lstrip(" "))
            items.append((min(indent // 2, 4), raw.strip()[1:].strip()))

        slides.append(
            {
                "title": title.strip(),
                "items": items,
                "diagram": diagram,
                "caption": caption,
                "notes": notes.strip("\n").rstrip(),
                "time": time,
            }
        )
    return meta, slides


# ------------------------------------------------------------------- text ---
def rich(para, text: str, size: int, color=INK, bold=False) -> None:
    """Write text into a paragraph, honouring **bold** and `code`."""
    for token in INLINE.split(text):
        if not token:
            continue
        run = para.add_run()
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT
        run.font.bold = bold
        if token.startswith("**") and token.endswith("**"):
            run.text = token[2:-2]
            run.font.bold = True
        elif token.startswith("`") and token.endswith("`"):
            run.text = token[1:-1]
            run.font.name = MONO
        else:
            run.text = token


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    return frame


def box_shape(slide, left, top, width, height, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.25)
    shape.shadow.inherit = False
    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.16)
    frame.margin_top = frame.margin_bottom = Inches(0.12)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.clear()
    return shape


def arrow(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def split_item(text: str) -> tuple[str, str]:
    label, _, detail = text.partition("|")
    return label.strip(), detail.strip()


def plain(text: str) -> str:
    """Text length as rendered, without markup characters."""
    return re.sub(r"[*`]", "", text)


def wrapped_lines(text: str, width_in: float, pt: float) -> int:
    """Approximate number of wrapped lines for text in a box of given width."""
    if not text.strip():
        return 0
    usable = max(width_in - 0.34, 0.4)
    char_w = pt / 185.0  # ~average glyph width in inches for Segoe UI
    per_line = max(int(usable / char_w), 8)
    return max(1, -(-len(plain(text)) // per_line))


def line_h(pt: float) -> float:
    return pt * 1.32 / 72.0


def fill_box(shape, label: str, detail: str, lsize: int, dsize: int, on_dark: bool):
    frame = shape.text_frame
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rich(p, label, lsize, WHITE if on_dark else ACCENT_DK, bold=True)
    if detail:
        d = frame.add_paragraph()
        d.alignment = PP_ALIGN.CENTER
        d.space_before = Pt(4)
        rich(d, detail, dsize, WHITE if on_dark else INK)


# --------------------------------------------------------------- diagrams ---
def _step_height(items, w_in: float, lpt: float, dpt: float) -> float:
    """Height that fits the tallest label+detail among items."""
    need = 0.0
    for _, text in items:
        label, detail = split_item(text)
        need = max(
            need,
            wrapped_lines(label, w_in, lpt) * line_h(lpt)
            + wrapped_lines(detail, w_in, dpt) * line_h(dpt),
        )
    return round(need + 0.34, 2)  # + vertical margins


def draw_ladder(slide, items, top_in: float, avail_in: float):
    """Ascending steps: progression / maturity levels."""
    n = max(len(items), 1)
    gap = 0.28
    w = (13.333 - 2 * 0.75 - gap * (n - 1)) / n
    rise = 0.72 if n > 1 else 0.0
    # Fill ~92% of the body region, but never shrink below what the text needs.
    h = max(_step_height(items, w, 18, 13), 0.92 * avail_in - rise * (n - 1), 1.3)
    group_h = h + rise * (n - 1)
    base_top = top_in + max((avail_in - group_h) / 2, 0) + rise * (n - 1)

    for i, (_, text) in enumerate(items):
        label, detail = split_item(text)
        left = 0.75 + i * (w + gap)
        top = base_top - rise * i
        tint = TINTS[min(int(i * (len(TINTS) - 1) / max(n - 1, 1)), len(TINTS) - 1)]
        on_dark = i >= n - 1 and n > 2
        shape = box_shape(
            slide, Inches(left), Inches(top), Inches(w), Inches(h), tint, ACCENT
        )
        fill_box(shape, label, detail, 18, 13, on_dark)

        if i < n - 1:
            arrow(
                slide,
                Inches(left + w + 0.02),
                Inches(top - rise + h / 2 - 0.11),
                Inches(gap - 0.04),
                Inches(0.22),
            )


def _draw_steps(slide, items, top_in, avail_in, lpt, dpt, loop=False):
    n = max(len(items), 1)
    gap = 0.40
    w = (13.333 - 2 * 0.75 - gap * (n - 1)) / n
    extra = 0.75 if loop else 0.0
    h = max(_step_height(items, w, lpt, dpt), 0.66 * (avail_in - extra), 1.2)
    top = top_in + max((avail_in - h - extra) / 2, 0)

    for i, (_, text) in enumerate(items):
        label, detail = split_item(text)
        left = 0.75 + i * (w + gap)
        shape = box_shape(
            slide, Inches(left), Inches(top), Inches(w), Inches(h), TINTS[0], ACCENT
        )
        fill_box(shape, label, detail, lpt, dpt, False)
        if i < n - 1:
            arrow(
                slide,
                Inches(left + w + 0.05),
                Inches(top + h / 2 - 0.13),
                Inches(gap - 0.10),
                Inches(0.26),
            )
    return top, h


def draw_flow(slide, items, top_in: float, avail_in: float):
    """Left-to-right process boxes with arrows."""
    _draw_steps(slide, items, top_in, avail_in, 16, 12.5)


def draw_cycle(slide, items, top_in: float, avail_in: float):
    """Boxes in a row with a return arrow underneath."""
    top, h = _draw_steps(slide, items, top_in, avail_in, 15, 11.5, loop=True)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.LEFT_ARROW,
        Inches(1.25),
        Inches(top + h + 0.28),
        Inches(11.33 - 0.5),
        Inches(0.28),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TINTS[1]
    bar.line.fill.background()
    bar.shadow.inherit = False
    p = bar.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rich(p, "repeat", 11, ACCENT_DK, bold=True)


def draw_compare(slide, items, top_in: float, avail_in: float):
    """Side-by-side columns: before/after, with/without."""
    n = max(len(items), 1)
    gap = 0.5
    w = (13.333 - 2 * 0.75 - gap * (n - 1)) / n
    head_h = 0.72
    bullet_pt = 16.0

    rows = 0.0
    for _, text in items:
        _, detail = split_item(text)
        lines = [x.strip() for x in detail.split(";") if x.strip()]
        rows = max(
            rows,
            sum(wrapped_lines(x, w, bullet_pt) for x in lines) + len(lines) * 0.45,
        )
    body_h = min(
        max(rows * line_h(bullet_pt) + 0.40, 0.72 * (avail_in - head_h), 1.1),
        avail_in - head_h - 0.16,
    )

    group_h = head_h + 0.08 + body_h
    top = top_in + max((avail_in - group_h) / 2, 0)

    for i, (_, text) in enumerate(items):
        label, detail = split_item(text)
        left = 0.75 + i * (w + gap)
        tint = ACCENT if i == n - 1 else RGBColor(0x86, 0x8B, 0x93)

        head = box_shape(
            slide, Inches(left), Inches(top), Inches(w), Inches(head_h), tint
        )
        hp = head.text_frame.paragraphs[0]
        hp.alignment = PP_ALIGN.CENTER
        rich(hp, label, 17, WHITE, bold=True)

        body = box_shape(
            slide,
            Inches(left),
            Inches(top + head_h + 0.08),
            Inches(w),
            Inches(body_h),
            RGBColor(0xF7, 0xF9, 0xFC),
            RGBColor(0xD5, 0xDF, 0xEA),
        )
        frame = body.text_frame
        frame.vertical_anchor = MSO_ANCHOR.TOP
        for j, line in enumerate(x.strip() for x in detail.split(";") if x.strip()):
            p = frame.paragraphs[0] if j == 0 else frame.add_paragraph()
            p.space_after = Pt(8)
            rich(p, "•  " + line, bullet_pt)


DIAGRAMS = {
    "ladder": draw_ladder,
    "flow": draw_flow,
    "compare": draw_compare,
    "cycle": draw_cycle,
}


# ---------------------------------------------------------------- builder ---
def add_title(slide, text: str, time: str) -> None:
    frame = textbox(slide, MARGIN, Inches(0.42), Inches(10.2), Inches(0.95))
    frame.vertical_anchor = MSO_ANCHOR.TOP
    rich(frame.paragraphs[0], text, 29, INK, bold=True)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.44), Inches(1.15), Inches(0.055)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    if time:
        tf = textbox(slide, Inches(11.35), Inches(0.5), Inches(1.3), Inches(0.4))
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        rich(tf.paragraphs[0], time, 12, ACCENT, bold=True)


def add_bullets(slide, items) -> None:
    frame = textbox(slide, MARGIN, BODY_TOP, BODY_W, BODY_H)
    size = 19 if len(items) <= 6 else 17
    for i, (level, text) in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.space_after = Pt(11 if level == 0 else 5)
        indent = "    " * level
        rich(p, f"{indent}{'–' if level else '•'}  {text}", size - level)


def build(meta: dict, slides: list[dict], out: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    for index, slide_spec in enumerate(slides):
        s = prs.slides.add_slide(blank)

        if index == 0:
            bar = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(0), Inches(2.55), SLIDE_W, Inches(0.06)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT
            bar.line.fill.background()
            bar.shadow.inherit = False

            frame = textbox(s, MARGIN, Inches(1.55), Inches(11.6), Inches(1.0))
            frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            rich(frame.paragraphs[0], slide_spec["title"], 38, INK, bold=True)

            sub = textbox(s, MARGIN, Inches(2.85), Inches(11.6), Inches(1.6))
            for j, line in enumerate(
                b
                for b in (
                    meta.get("subtitle", ""),
                    meta.get("presenter", ""),
                    meta.get("ppt_date_time", ""),
                )
                if b
            ):
                p = sub.paragraphs[0] if j == 0 else sub.add_paragraph()
                p.space_after = Pt(6)
                rich(p, line, 17, MUTED)
        else:
            add_title(s, slide_spec["title"], slide_spec["time"])

            caption = slide_spec["caption"]
            top_in = BODY_TOP / 914400
            avail_in = BODY_H / 914400 - (0.62 if caption else 0.0)

            drawer = DIAGRAMS.get(slide_spec["diagram"])
            if drawer and slide_spec["items"]:
                drawer(s, slide_spec["items"], top_in, avail_in)
            else:
                add_bullets(s, slide_spec["items"])

            if caption:
                cf = textbox(
                    s,
                    MARGIN,
                    Inches(top_in + avail_in + 0.12),
                    BODY_W,
                    Inches(0.45),
                )
                cf.paragraphs[0].alignment = PP_ALIGN.CENTER
                rich(cf.paragraphs[0], caption, 16, ACCENT_DK, bold=True)

        if slide_spec["notes"]:
            s.notes_slide.notes_text_frame.text = slide_spec["notes"]

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    ap.add_argument("deck", type=Path, help="path to the .deck.md file")
    ap.add_argument("-o", "--out", type=Path, help="output .pptx path")
    args = ap.parse_args()

    meta, slides = parse_deck(args.deck.read_text(encoding="utf-8"))
    if not slides:
        sys.exit(f"No '## ' slides found in {args.deck}")

    out = args.out or args.deck.with_suffix("").with_suffix(".pptx")
    try:
        build(meta, slides, out)
    except PermissionError:
        alt = out.with_name(out.stem + ".new.pptx")
        build(meta, slides, alt)
        sys.exit(
            f"'{out.name}' is open in PowerPoint, so it could not be replaced.\n"
            f"Wrote {alt} instead — close PowerPoint and either rename it, or\n"
            f"re-run this command to overwrite the original."
        )

    noted = sum(1 for s in slides if s["notes"])
    diagrams = sum(1 for s in slides if s["diagram"] in DIAGRAMS)
    print(
        f"Wrote {out}  —  {len(slides)} slides, "
        f"{noted} with speaker notes, {diagrams} diagrams"
    )


if __name__ == "__main__":
    main()
